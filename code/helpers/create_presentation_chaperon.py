"""
Create an engaging PowerPoint presentation explaining the Bayesian Hiring Model
"""

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
import os

# Set style for plots
sns.set_style("whitegrid")
plt.rcParams['figure.facecolor'] = 'white'
plt.rcParams['axes.facecolor'] = 'white'


def create_probability_matrix():
    """Create the probability matrix visualization"""
    try:
        from batch_simulation_chaperon import BayesianHiringModel
    except ImportError:
        # Fallback: recreate the model here
        import sys
        sys.path.insert(0, '.')
        from batch_simulation_chaperon import BayesianHiringModel
    
    model = BayesianHiringModel(
        p=0.5,
        q_good=0.7,
        q_bad=0.3,
        lambda_good=5.0,
        lambda_bad=2.0
    )
    return model.compute_probability_matrix(n_max=10)


def create_heatmap_figure(matrix, figsize=(10, 5)):
    """Create a heatmap visualization"""
    fig, ax = plt.subplots(figsize=figsize)
    sns.heatmap(matrix, annot=True, fmt='.2f', cmap='RdYlGn', 
               vmin=0, vmax=1, cbar_kws={'label': 'Probability'}, ax=ax)
    ax.set_title('Probability of Being Good', fontsize=14, fontweight='bold')
    ax.set_xlabel('Number of Papers (N)', fontsize=12)
    ax.set_ylabel('Network Effect (G)', fontsize=12)
    plt.tight_layout()
    return fig


def create_prior_visualization():
    """Visualize the prior distribution"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4))
    
    # Prior distribution
    categories = ['Good', 'Bad']
    probs = [0.5, 0.5]
    colors = ['#2ecc71', '#e74c3c']
    ax1.bar(categories, probs, color=colors, alpha=0.7, edgecolor='black', linewidth=2)
    ax1.set_ylim(0, 1)
    ax1.set_ylabel('Probability', fontsize=12)
    ax1.set_title('Prior: Before Observing Signals', fontsize=14, fontweight='bold')
    ax1.grid(axis='y', alpha=0.3)
    for i, (cat, prob) in enumerate(zip(categories, probs)):
        ax1.text(i, prob + 0.05, f'{prob:.1%}', ha='center', fontsize=14, fontweight='bold')
    
    # Network effect likelihood
    network_data = {
        'Good': {'Yes': 0.7, 'No': 0.3},
        'Bad': {'Yes': 0.3, 'No': 0.7}
    }
    x = np.arange(len(['Yes', 'No']))
    width = 0.35
    ax2.bar(x - width/2, [network_data['Good']['Yes'], network_data['Good']['No']], 
           width, label='Good', color='#2ecc71', alpha=0.7, edgecolor='black')
    ax2.bar(x + width/2, [network_data['Bad']['Yes'], network_data['Bad']['No']], 
           width, label='Bad', color='#e74c3c', alpha=0.7, edgecolor='black')
    ax2.set_xlabel('Network Effect', fontsize=12)
    ax2.set_ylabel('Probability', fontsize=12)
    ax2.set_title('Network Signal: Pr(G | F)', fontsize=14, fontweight='bold')
    ax2.set_xticks(x)
    ax2.set_xticklabels(['Yes', 'No'])
    ax2.legend()
    ax2.grid(axis='y', alpha=0.3)
    ax2.set_ylim(0, 1)
    
    plt.tight_layout()
    return fig


def create_poisson_visualization():
    """Visualize Poisson distributions for paper counts"""
    fig, ax = plt.subplots(figsize=(10, 5))
    
    n_range = np.arange(0, 12)
    lambda_good = 5.0
    lambda_bad = 2.0
    
    from scipy.stats import poisson
    pmf_good = poisson.pmf(n_range, lambda_good)
    pmf_bad = poisson.pmf(n_range, lambda_bad)
    
    ax.bar(n_range - 0.2, pmf_good, width=0.4, label='Good Candidates', 
          color='#2ecc71', alpha=0.7, edgecolor='black')
    ax.bar(n_range + 0.2, pmf_bad, width=0.4, label='Bad Candidates', 
          color='#e74c3c', alpha=0.7, edgecolor='black')
    
    ax.set_xlabel('Number of Papers (N)', fontsize=12)
    ax.set_ylabel('Probability', fontsize=12)
    ax.set_title('Publication Count Signal: Pr(N | F)', fontsize=14, fontweight='bold')
    ax.legend(fontsize=11)
    ax.grid(axis='y', alpha=0.3)
    ax.set_xticks(n_range)
    
    plt.tight_layout()
    return fig


def create_bayes_visualization():
    """Visualize Bayes' theorem concept"""
    fig, ax = plt.subplots(figsize=(10, 6))
    ax.axis('off')
    
    # Create a visual representation of Bayes' theorem
    text = r"""
    Bayes' Theorem: Updating Beliefs
    
    $\Pr(F = \text{good} | G, N) = \frac{\Pr(G, N | F = \text{good}) \cdot \Pr(F = \text{good})}{\Pr(G, N)}$
    
    Prior × Likelihood = Posterior
    """
    
    ax.text(0.5, 0.5, text, fontsize=16, ha='center', va='center',
           bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
    
    return fig


def create_example_scenarios():
    """Create example candidate scenarios"""
    try:
        from batch_simulation_chaperon import BayesianHiringModel
    except ImportError:
        import sys
        sys.path.insert(0, '.')
        from batch_simulation_chaperon import BayesianHiringModel
    
    model = BayesianHiringModel()
    
    scenarios = [
        ("Candidate A", True, 8, "Strong"),
        ("Candidate B", True, 2, "Moderate"),
        ("Candidate C", False, 6, "Moderate"),
        ("Candidate D", False, 1, "Weak")
    ]
    
    fig, axes = plt.subplots(2, 2, figsize=(12, 10))
    axes = axes.flatten()
    
    for idx, (name, has_network, n_papers, strength) in enumerate(scenarios):
        ax = axes[idx]
        posterior_good, posterior_bad = model.posterior_probability(has_network, n_papers)
        
        # Create bar chart
        categories = ['Good', 'Bad']
        probs = [posterior_good, posterior_bad]
        colors = ['#2ecc71' if posterior_good > 0.5 else '#95a5a6', 
                 '#e74c3c' if posterior_bad > 0.5 else '#95a5a6']
        
        bars = ax.bar(categories, probs, color=colors, alpha=0.7, edgecolor='black', linewidth=2)
        ax.set_ylim(0, 1)
        ax.set_ylabel('Probability', fontsize=11)
        ax.set_title(f'{name}\nNetwork: {"Yes" if has_network else "No"}, Papers: {n_papers}', 
                    fontsize=12, fontweight='bold')
        ax.grid(axis='y', alpha=0.3)
        
        for bar, prob in zip(bars, probs):
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 0.02,
                   f'{prob:.1%}', ha='center', fontsize=12, fontweight='bold')
        
        # Add decision
        decision = "HIRE" if posterior_good > 0.5 else "REJECT"
        ax.text(0.5, 0.9, decision, transform=ax.transAxes, 
               ha='center', fontsize=14, fontweight='bold',
               bbox=dict(boxstyle='round', facecolor='yellow', alpha=0.7))
    
    plt.tight_layout()
    return fig


def save_figure_to_image(fig, filename):
    """Save matplotlib figure to temporary file"""
    fig.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return filename


def add_slide_with_image(prs, title, content_text, image_path=None, layout_idx=1):
    """Add a slide with optional image"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    if content_text:
        if layout_idx == 0:  # Title slide
            content = slide.placeholders[1]
        else:
            # Add text box
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(4)
            content = slide.shapes.add_textbox(left, top, width, height).text_frame
        
        content.text = content_text
        content.paragraphs[0].font.size = Pt(18)
        content.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    
    # Image
    if image_path and os.path.exists(image_path):
        left = Inches(1) if layout_idx == 1 else Inches(2)
        top = Inches(2.5) if layout_idx == 1 else Inches(3)
        slide.shapes.add_picture(image_path, left, top, width=Inches(8), height=Inches(4.5))
    
    return slide


def create_presentation(output_dir=None):
    """
    Create the full PowerPoint presentation
    
    Parameters:
    -----------
    output_dir : str, optional
        Directory where to save the presentation file. If None, saves in current directory.
        If directory doesn't exist, it will be created.
    
    Returns:
    --------
    tuple : (Presentation object, output_file_path)
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    primary_color = RGBColor(0, 51, 102)
    accent_color = RGBColor(46, 204, 113)
    danger_color = RGBColor(231, 76, 60)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Bayesian Hiring Model"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    subtitle.text = "Using Network Effects and Publication Count\nTo Make Better Hiring Decisions"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Hiring Challenge"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🎯 You need to hire the best candidates"
    p = tf.add_paragraph()
    p.text = "❓ But you can't directly observe their true quality"
    p = tf.add_paragraph()
    p.text = "📊 You only see signals:"
    p = tf.add_paragraph()
    p.text = "   • Network effect (G): Did they work with top researchers?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Publication count (N): How many papers did they publish?"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "💡 Solution: Use Bayesian inference to combine signals!"
    p.font.size = Pt(20)
    p.font.bold = True
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)
    
    # Slide 3: The Model Overview
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "What We're Modeling"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Create visual diagram
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    
    # Add text boxes for the model components
    # F (Fit) - hidden
    fit_box = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(1.5), Inches(0.8))
    fit_box.fill.solid()
    fit_box.fill.fore_color.rgb = RGBColor(200, 200, 200)
    fit_box.text = "F (Fit)\nHidden"
    fit_box.text_frame.paragraphs[0].font.size = Pt(14)
    fit_box.text_frame.paragraphs[0].font.bold = True
    
    # G (Network) - observed
    network_box = slide.shapes.add_shape(1, Inches(3.5), Inches(2.5), Inches(1.5), Inches(0.8))
    network_box.fill.solid()
    network_box.fill.fore_color.rgb = accent_color
    network_box.text = "G (Network)\nObserved"
    network_box.text_frame.paragraphs[0].font.size = Pt(14)
    network_box.text_frame.paragraphs[0].font.bold = True
    
    # N (Papers) - observed
    papers_box = slide.shapes.add_shape(1, Inches(6), Inches(2.5), Inches(1.5), Inches(0.8))
    papers_box.fill.solid()
    papers_box.fill.fore_color.rgb = accent_color
    papers_box.text = "N (Papers)\nObserved"
    papers_box.text_frame.paragraphs[0].font.size = Pt(14)
    papers_box.text_frame.paragraphs[0].font.bold = True
    
    # Arrows (using lines with arrowheads)
    from pptx.enum.shapes import MSO_SHAPE
    # Create simple lines as arrows
    line1 = slide.shapes.add_connector(1, Inches(2.5), Inches(2.9), Inches(3.5), Inches(2.9))
    line2 = slide.shapes.add_connector(1, Inches(2.5), Inches(2.9), Inches(6), Inches(2.9))
    
    # Goal text
    goal_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    goal_frame = goal_box.text_frame
    goal_frame.text = "Goal: Compute Pr(F = good | G, N)"
    goal_frame.paragraphs[0].font.size = Pt(20)
    goal_frame.paragraphs[0].font.bold = True
    goal_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Slide 4: Step 1 - Prior
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 1: Start with Prior Belief"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Before seeing any signals, we have an initial belief:"
    p = tf.add_paragraph()
    p.text = "Pr(F = good) = p = 0.5"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p = tf.add_paragraph()
    p.text = "This means: 50% of candidates are good (on average)"
    p.font.size = Pt(18)
    
    # Add visualization
    fig = create_prior_visualization()
    img_path = 'temp_prior.png'
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(1), Inches(3.5), width=Inches(8), height=Inches(3))
    
    # Slide 5: Step 2 - Network Signal
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 2: Network Effect Signal"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Network effect tells us about candidate quality:"
    p = tf.add_paragraph()
    p.text = "• Good candidates: 70% have network (q_good = 0.7)"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "• Bad candidates: 30% have network (q_bad = 0.3)"
    p.font.size = Pt(18)
    p.font.color.rgb = danger_color
    p = tf.add_paragraph()
    p.text = "Key insight: Network effect is a stronger signal for good candidates!"
    p.font.size = Pt(18)
    p.font.bold = True
    
    # Slide 6: Step 3 - Publication Signal
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 3: Publication Count Signal"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Publication count follows Poisson distribution:"
    p = tf.add_paragraph()
    p.text = "• Good candidates: Average 5 papers (λ_good = 5.0)"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "• Bad candidates: Average 2 papers (λ_bad = 2.0)"
    p.font.size = Pt(18)
    p.font.color.rgb = danger_color
    
    # Add visualization
    fig = create_poisson_visualization()
    img_path = 'temp_poisson.png'
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(3), width=Inches(9), height=Inches(4))
    
    # Slide 7: Bayes' Theorem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 4: Combine Signals with Bayes' Theorem"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Bayes' Theorem combines prior and signals:"
    p = tf.add_paragraph()
    p.text = "Pr(F = good | G, N) ="
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = "    Pr(G, N | good) × Pr(good)"
    p.font.size = Pt(18)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "    ─────────────────────────────"
    p.font.size = Pt(18)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "         Pr(G, N)"
    p.font.size = Pt(18)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Think of it as: Prior × Evidence = Updated Belief"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    # Slide 8: The Probability Matrix
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Result: Probability Matrix"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "For each combination of signals, we compute the probability:"
    p = tf.add_paragraph()
    p.text = "Rows: Network effect (Yes/No)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "Columns: Number of papers (1-10)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "Values: Pr(F = good | G, N)"
    p.font.size = Pt(18)
    p.font.bold = True
    
    # Add heatmap
    matrix = create_probability_matrix()
    fig = create_heatmap_figure(matrix)
    img_path = 'temp_heatmap.png'
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(3.5))
    
    # Slide 9: Example Scenarios
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Real Examples: Four Candidates"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    # Add visualization
    fig = create_example_scenarios()
    img_path = 'temp_examples.png'
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5.5))
    
    # Slide 10: Key Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Insights"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🎯 Network effect is powerful:"
    p = tf.add_paragraph()
    p.text = "   Candidates with network + few papers can still be good"
    p.level = 1
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "📊 Many papers help, but network matters more:"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "   No network + many papers = moderate probability"
    p.level = 1
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "⚠️ Weakest signal: No network + few papers"
    p.font.size = Pt(18)
    p.font.color.rgb = danger_color
    p = tf.add_paragraph()
    p.text = "✅ Strongest signal: Network + many papers"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    
    # Slide 11: Decision Rule
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "How to Use: Decision Rule"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Simple decision rule:"
    p = tf.add_paragraph()
    p.text = "IF Pr(F = good | G, N) > 0.5"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "   → HIRE"
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "ELSE"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = danger_color
    p = tf.add_paragraph()
    p.text = "   → REJECT"
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Slide 12: Why This Matters
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why This Model Matters"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "✅ Combines multiple signals systematically"
    p = tf.add_paragraph()
    p.text = "✅ Quantifies uncertainty (probabilities, not just yes/no)"
    p = tf.add_paragraph()
    p.text = "✅ Updates beliefs based on evidence"
    p = tf.add_paragraph()
    p.text = "✅ Can be calibrated with real hiring data"
    p = tf.add_paragraph()
    p.text = "✅ Transparent and interpretable"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(10)
    
    # Slide 13: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1️⃣ Start with prior belief about candidate quality"
    p = tf.add_paragraph()
    p.text = "2️⃣ Observe signals: Network effect (G) and Papers (N)"
    p = tf.add_paragraph()
    p.text = "3️⃣ Use Bayes' theorem to update beliefs"
    p = tf.add_paragraph()
    p.text = "4️⃣ Get posterior probability: Pr(F = good | G, N)"
    p = tf.add_paragraph()
    p.text = "5️⃣ Make hiring decision based on probability"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(22)
        paragraph.space_after = Pt(12)
        paragraph.font.bold = True
    
    # Slide 14: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Thank You!"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Questions?"
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Clean up temporary files
    temp_files = ['temp_prior.png', 'temp_poisson.png', 'temp_heatmap.png', 'temp_examples.png']
    for f in temp_files:
        if os.path.exists(f):
            os.remove(f)
    
    # Determine output file path
    output_filename = "Chaperon_Based_Hiring_Model.pptx"
    if output_dir:
        # Create directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, output_filename)
    else:
        output_path = output_filename
    
    # Save the presentation
    prs.save(output_path)
    
    return prs, output_path


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(
        description='Create PowerPoint presentation for Bayesian Hiring Model'
    )
    parser.add_argument(
        '--dir', '-d',
        type=str,
        default=None,
        help='Directory where to save the presentation file (default: current directory)'
    )
    
    args = parser.parse_args()
    
    print("Creating PowerPoint presentation...")
    prs, output_file = create_presentation(output_dir=args.dir)
    print(f"✅ Presentation saved as '{output_file}'")
    print(f"📊 Total slides: {len(prs.slides)}")
