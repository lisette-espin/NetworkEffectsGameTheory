"""
Create an engaging PowerPoint presentation explaining the Referral-Based (FOF) Bayesian Hiring Model
"""

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from io import BytesIO
import os

# Set style for plots
sns.set_style("whitegrid")
plt.rcParams['figure.facecolor'] = 'white'
plt.rcParams['axes.facecolor'] = 'white'


def create_probability_matrix():
    """Create the probability matrix visualization"""
    try:
        from batch_simulation_fof import ReferralBasedHiringModel
    except ImportError:
        import sys
        sys.path.append('.')
        sys.path.append('../')
        sys.path.append('code')
        from batch_simulation_fof import ReferralBasedHiringModel
    
    model = ReferralBasedHiringModel(
        p=0.5,
        lambda_good=5.0,
        lambda_bad=2.0,
        referral_accuracy=0.8,
        referral_bias=0.1
    )
    return model.compute_probability_matrix(n_max=10)


def create_heatmap_figure(matrix, figsize=(12, 8)):
    """Create a heatmap visualization"""
    fig, ax = plt.subplots(figsize=figsize)
    sns.heatmap(matrix, annot=True, fmt='.2f', cmap='RdYlGn', 
               vmin=0, vmax=1, cbar_kws={'label': 'Probability'}, ax=ax)
    ax.set_title('Probability of Being Good Given Network Effect (G) and Paper Count (N)', 
                 fontsize=14, fontweight='bold')
    ax.set_xlabel('Number of Papers (N)', fontsize=12)
    ax.set_ylabel('Network Effect (G) = Referral Quality', fontsize=12)
    plt.tight_layout()
    return fig


def create_prior_visualization():
    """Visualize the prior distribution"""
    fig, ax = plt.subplots(1, 1, figsize=(8, 5))
    
    # Prior distribution
    categories = ['Good', 'Bad']
    probs = [0.5, 0.5]
    colors = ['#2ecc71', '#e74c3c']
    ax.bar(categories, probs, color=colors, alpha=0.7, edgecolor='black', linewidth=2)
    ax.set_ylim(0, 1)
    ax.set_ylabel('Probability', fontsize=12)
    ax.set_title('Prior: Before Observing Signals', fontsize=14, fontweight='bold')
    ax.grid(axis='y', alpha=0.3)
    for i, (cat, prob) in enumerate(zip(categories, probs)):
        ax.text(i, prob + 0.05, f'{prob:.1%}', ha='center', fontsize=14, fontweight='bold')
    
    plt.tight_layout()
    return fig


def create_referral_distribution_visualization():
    """Visualize referral quality distributions"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5))
    
    # Good candidates
    qualities = ['Excellent', 'Good', 'Average', 'Bad', 'Very Bad']
    probs_good = [0.33, 0.25, 0.18, 0.12, 0.12]
    colors_good = ['#2ecc71', '#27ae60', '#f39c12', '#e67e22', '#c0392b']
    
    ax1.bar(qualities, probs_good, color=colors_good, alpha=0.7, edgecolor='black', linewidth=1.5)
    ax1.set_ylim(0, 0.4)
    ax1.set_ylabel('Probability', fontsize=11)
    ax1.set_title('Referral Quality Given F=Good', fontsize=13, fontweight='bold')
    ax1.grid(axis='y', alpha=0.3)
    ax1.tick_params(axis='x', rotation=45)
    for i, (q, p) in enumerate(zip(qualities, probs_good)):
        ax1.text(i, p + 0.02, f'{p:.2f}', ha='center', fontsize=10, fontweight='bold')
    
    # Bad candidates
    probs_bad = [0.05, 0.10, 0.15, 0.30, 0.40]
    colors_bad = ['#2ecc71', '#27ae60', '#f39c12', '#e67e22', '#c0392b']
    
    ax2.bar(qualities, probs_bad, color=colors_bad, alpha=0.7, edgecolor='black', linewidth=1.5)
    ax2.set_ylim(0, 0.4)
    ax2.set_ylabel('Probability', fontsize=11)
    ax2.set_title('Referral Quality Given F=Bad', fontsize=13, fontweight='bold')
    ax2.grid(axis='y', alpha=0.3)
    ax2.tick_params(axis='x', rotation=45)
    for i, (q, p) in enumerate(zip(qualities, probs_bad)):
        ax2.text(i, p + 0.02, f'{p:.2f}', ha='center', fontsize=10, fontweight='bold')
    
    plt.tight_layout()
    return fig


def create_poisson_visualization():
    """Visualize Poisson distributions for paper counts"""
    fig, ax = plt.subplots(figsize=(10, 5))
    
    n_range = np.arange(0, 12)
    lambda_good = 5.0
    lambda_bad = 2.0
    
    from scipy.stats import poisson
    pmf_good = poisson.pmf(n_range, lambda_good)
    pmf_bad = poisson.pmf(n_range, lambda_bad)
    
    ax.bar(n_range - 0.2, pmf_good, width=0.4, label='Good Candidates', 
          color='#2ecc71', alpha=0.7, edgecolor='black')
    ax.bar(n_range + 0.2, pmf_bad, width=0.4, label='Bad Candidates', 
          color='#e74c3c', alpha=0.7, edgecolor='black')
    
    ax.set_xlabel('Number of Papers (N)', fontsize=12)
    ax.set_ylabel('Probability', fontsize=12)
    ax.set_title('Publication Count Signal: Pr(N | F)', fontsize=14, fontweight='bold')
    ax.legend(fontsize=11)
    ax.grid(axis='y', alpha=0.3)
    ax.set_xticks(n_range)
    
    plt.tight_layout()
    return fig


def create_trust_selection_visualization():
    """Visualize trust-based referral selection scenarios"""
    fig, axes = plt.subplots(2, 2, figsize=(14, 10))
    axes = axes.flatten()
    
    scenarios = [
        ("Trust R1 Only", "R1: Excellent\nR2: Very Bad\n→ Use R1", "#2ecc71"),
        ("Trust R2 Only", "R1: Excellent\nR2: Very Bad\n→ Use R2", "#e74c3c"),
        ("Trust Both\n(Conservative)", "R1: Good\nR2: Bad\n→ Use R2 (worst)", "#f39c12"),
        ("Trust Neither", "R1: Excellent\nR2: Good\n→ Use R1 (best)", "#3498db")
    ]
    
    for idx, (title, description, color) in enumerate(scenarios):
        ax = axes[idx]
        ax.axis('off')
        ax.text(0.5, 0.7, title, transform=ax.transAxes, ha='center', 
               fontsize=14, fontweight='bold', color=color)
        ax.text(0.5, 0.4, description, transform=ax.transAxes, ha='center', 
               fontsize=11, bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
    
    plt.tight_layout()
    return fig


def create_example_scenarios():
    """Create example candidate scenarios"""
    try:
        from batch_simulation_fof import ReferralBasedHiringModel
    except ImportError:
        import sys
        sys.path.append('.')
        sys.path.append('../')
        sys.path.append('code')
        from batch_simulation_fof import ReferralBasedHiringModel
    
    model = ReferralBasedHiringModel()
    
    scenarios = [
        ("Candidate A", "excellent", 8, "Strong"),
        ("Candidate B", "good", 5, "Moderate"),
        ("Candidate C", "average", 3, "Weak"),
        ("Candidate D", "very bad", 1, "Very Weak")
    ]
    
    fig, axes = plt.subplots(2, 2, figsize=(12, 10))
    axes = axes.flatten()
    
    for idx, (name, network_level, n_papers, strength) in enumerate(scenarios):
        ax = axes[idx]
        posterior_good, posterior_bad = model.posterior_probability(network_level, n_papers)
        
        # Create bar chart
        categories = ['Good', 'Bad']
        probs = [posterior_good, posterior_bad]
        colors = ['#2ecc71' if posterior_good > 0.5 else '#95a5a6', 
                 '#e74c3c' if posterior_bad > 0.5 else '#95a5a6']
        
        bars = ax.bar(categories, probs, color=colors, alpha=0.7, edgecolor='black', linewidth=2)
        ax.set_ylim(0, 1)
        ax.set_ylabel('Probability', fontsize=11)
        ax.set_title(f'{name}\nG={network_level.title()}, N={n_papers}', 
                    fontsize=12, fontweight='bold')
        ax.grid(axis='y', alpha=0.3)
        
        for bar, prob in zip(bars, probs):
            height = bar.get_height()
            ax.text(bar.get_x() + bar.get_width()/2., height + 0.02,
                   f'{prob:.1%}', ha='center', fontsize=12, fontweight='bold')
        
        # Add decision
        decision = "HIRE" if posterior_good > 0.5 else "REJECT"
        ax.text(0.5, 0.9, decision, transform=ax.transAxes, 
               ha='center', fontsize=14, fontweight='bold',
               bbox=dict(boxstyle='round', facecolor='yellow', alpha=0.7))
    
    plt.tight_layout()
    return fig


def create_payoff_visualization():
    """Visualize payoff structure"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 5))
    
    # Payoff structure
    B = 10.0
    b = -2.0  # Note: user changed default to -2.0
    w = 5.0
    
    # Employer payoffs
    categories = ['Hire Good', 'Hire Bad', 'Reject']
    employer_payoffs = [B, b, 0]
    colors_emp = ['#2ecc71', '#e74c3c', '#95a5a6']
    
    bars1 = ax1.bar(categories, employer_payoffs, color=colors_emp, alpha=0.7, 
                   edgecolor='black', linewidth=2)
    ax1.set_ylabel('Payoff', fontsize=12)
    ax1.set_title('Employer Payoffs', fontsize=14, fontweight='bold')
    ax1.grid(axis='y', alpha=0.3)
    ax1.axhline(y=0, color='black', linestyle='--', linewidth=1)
    
    for bar, payoff in zip(bars1, employer_payoffs):
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height + (0.5 if height >= 0 else -1.5),
                f'{payoff:.1f}', ha='center', fontsize=12, fontweight='bold')
    
    # Applicant payoffs
    applicant_payoffs = [w, w, 0]
    bars2 = ax2.bar(categories, applicant_payoffs, color=['#3498db', '#3498db', '#95a5a6'], 
                   alpha=0.7, edgecolor='black', linewidth=2)
    ax2.set_ylabel('Payoff', fontsize=12)
    ax2.set_title('Applicant Payoffs', fontsize=14, fontweight='bold')
    ax2.grid(axis='y', alpha=0.3)
    ax2.axhline(y=0, color='black', linestyle='--', linewidth=1)
    
    for bar, payoff in zip(bars2, applicant_payoffs):
        height = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width()/2., height + 0.3,
                f'{payoff:.1f}', ha='center', fontsize=12, fontweight='bold')
    
    plt.tight_layout()
    return fig


def create_expected_payoff_comparison():
    """Compare expected payoffs for different scenarios"""
    try:
        from batch_simulation_fof import ReferralBasedHiringModel
    except ImportError:
        import sys
        sys.path.append('.')
        sys.path.append('../')
        sys.path.append('code')
        from batch_simulation_fof import ReferralBasedHiringModel
    
    model = ReferralBasedHiringModel(B=10.0, b=-2.0, w=5.0)
    
    scenarios = [
        ("Excellent + Many Papers", "excellent", 8),
        ("Good + Medium Papers", "good", 5),
        ("Average + Few Papers", "average", 3),
        ("Very Bad + Few Papers", "very bad", 1)
    ]
    
    fig, ax = plt.subplots(figsize=(12, 6))
    
    scenario_names = []
    expected_payoffs = []
    posterior_probs = []
    
    for name, network_level, n_papers in scenarios:
        posterior_good, _ = model.posterior_probability(network_level, n_papers)
        expected_payoff = model.expected_payoff(network_level, n_papers)
        
        scenario_names.append(name)
        expected_payoffs.append(expected_payoff)
        posterior_probs.append(posterior_good)
    
    x = np.arange(len(scenario_names))
    width = 0.35
    
    bars1 = ax.bar(x - width/2, expected_payoffs, width, label='Expected Payoff', 
                   color='#3498db', alpha=0.7, edgecolor='black')
    bars2 = ax.bar(x + width/2, [p * 12 for p in posterior_probs], width, 
                   label='Posterior × 12 (scaled)', color='#e74c3c', alpha=0.7, edgecolor='black')
    
    ax.set_ylabel('Value', fontsize=12)
    ax.set_title('Expected Payoff vs. Posterior Probability', fontsize=14, fontweight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels(scenario_names, rotation=15, ha='right')
    ax.legend(fontsize=11)
    ax.grid(axis='y', alpha=0.3)
    ax.axhline(y=0, color='black', linestyle='--', linewidth=1)
    
    # Add value labels
    for bar in bars1:
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height + (0.5 if height >= 0 else -1.5),
               f'{height:.2f}', ha='center', fontsize=10, fontweight='bold')
    
    plt.tight_layout()
    return fig


def save_figure_to_image(fig, filename):
    """Save matplotlib figure to temporary file"""
    fig.savefig(filename, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return filename


def create_presentation(output_dir=None):
    """
    Create the full PowerPoint presentation
    
    Parameters:
    -----------
    output_dir : str, optional
        Directory where to save the presentation file. If None, saves in current directory.
        If directory doesn't exist, it will be created.
    
    Returns:
    --------
    Presentation : The created PowerPoint presentation object
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    primary_color = RGBColor(0, 51, 102)
    accent_color = RGBColor(46, 204, 113)
    danger_color = RGBColor(231, 76, 60)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Referral-Based Hiring Model"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    subtitle.text = "Trust, Referrals, and Network Effects\nA Bayesian Approach to Hiring Decisions"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Hiring Challenge with Referrals"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🎯 You receive TWO referrals for each candidate"
    p = tf.add_paragraph()
    p.text = "❓ But you don't trust all referrers equally"
    p = tf.add_paragraph()
    p.text = "💡 Key insight: Trust matters more than referral quality!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p = tf.add_paragraph()
    p.text = "📊 Example:"
    p = tf.add_paragraph()
    p.text = "   • Referral 1: 'Excellent' (from stranger)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   • Referral 2: 'Very Bad' (from trusted colleague)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "   → You use Referral 2 because you trust that source!"
    p.level = 1
    p.font.bold = True
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)
    
    # Slide 3: The Model Overview
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "What We're Modeling"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Add text boxes for the model components
    # F (Fit) - hidden
    fit_box = slide.shapes.add_shape(1, Inches(1), Inches(2), Inches(1.5), Inches(0.8))
    fit_box.fill.solid()
    fit_box.fill.fore_color.rgb = RGBColor(200, 200, 200)
    fit_box.text = "F (Fit)\nHidden"
    fit_box.text_frame.paragraphs[0].font.size = Pt(14)
    fit_box.text_frame.paragraphs[0].font.bold = True
    
    # R1, R2 (Referrals) - observed
    r1_box = slide.shapes.add_shape(1, Inches(3), Inches(2), Inches(1.5), Inches(0.8))
    r1_box.fill.solid()
    r1_box.fill.fore_color.rgb = accent_color
    r1_box.text = "R1, R2\nReferrals"
    r1_box.text_frame.paragraphs[0].font.size = Pt(14)
    r1_box.text_frame.paragraphs[0].font.bold = True
    
    # T1, T2 (Trust) - known
    t_box = slide.shapes.add_shape(1, Inches(5), Inches(2), Inches(1.5), Inches(0.8))
    t_box.fill.solid()
    t_box.fill.fore_color.rgb = RGBColor(52, 152, 219)
    t_box.text = "T1, T2\nTrust"
    t_box.text_frame.paragraphs[0].font.size = Pt(14)
    t_box.text_frame.paragraphs[0].font.bold = True
    
    # G (Network) - derived
    g_box = slide.shapes.add_shape(1, Inches(7), Inches(2), Inches(1.5), Inches(0.8))
    g_box.fill.solid()
    g_box.fill.fore_color.rgb = RGBColor(155, 89, 182)
    g_box.text = "G (Network)\n= Selected R"
    g_box.text_frame.paragraphs[0].font.size = Pt(14)
    g_box.text_frame.paragraphs[0].font.bold = True
    
    # N (Papers) - observed
    n_box = slide.shapes.add_shape(1, Inches(3), Inches(3.5), Inches(1.5), Inches(0.8))
    n_box.fill.solid()
    n_box.fill.fore_color.rgb = accent_color
    n_box.text = "N (Papers)\nObserved"
    n_box.text_frame.paragraphs[0].font.size = Pt(14)
    n_box.text_frame.paragraphs[0].font.bold = True
    
    # Goal text
    goal_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    goal_frame = goal_box.text_frame
    goal_frame.text = "Goal: Compute Pr(F = good | G, N)"
    goal_frame.paragraphs[0].font.size = Pt(20)
    goal_frame.paragraphs[0].font.bold = True
    goal_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Slide 4: Step 1 - Prior
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 1: Start with Prior Belief"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Before seeing any signals, we have an initial belief:"
    p = tf.add_paragraph()
    p.text = "Pr(F = good) = p = 0.5"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p = tf.add_paragraph()
    p.text = "This means: 50% of candidates are good (on average)"
    p.font.size = Pt(18)
    
    # Add visualization
    fig = create_prior_visualization()
    img_path = 'temp_prior_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(1), Inches(3.5), width=Inches(8), height=Inches(3))
    
    # Slide 5: Step 2 - Referral Quality
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 2: Referral Quality Signal"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Referrals come in 5 quality levels:"
    p = tf.add_paragraph()
    p.text = "Excellent → Good → Average → Bad → Very Bad"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Quality depends on candidate's true fit:"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Good candidates: More likely to get excellent/good referrals"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "• Bad candidates: More likely to get bad/very bad referrals"
    p.font.size = Pt(18)
    p.font.color.rgb = danger_color
    
    # Add visualization
    fig = create_referral_distribution_visualization()
    img_path = 'temp_referral_dist_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(3), width=Inches(9), height=Inches(4))
    
    # Slide 6: Step 3 - Trust-Based Selection
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 3: Trust-Based Referral Selection"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Decision Rule:"
    p = tf.add_paragraph()
    p.text = "1. Trust only R1 → Use R1 (regardless of R2)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "2. Trust only R2 → Use R2 (regardless of R1)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "3. Trust both → Use WORST referral (conservative!)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 126, 34)
    p = tf.add_paragraph()
    p.text = "4. Trust neither → Use BEST referral"
    p.font.size = Pt(18)
    
    # Add visualization
    fig = create_trust_selection_visualization()
    img_path = 'temp_trust_selection_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(3.5))
    
    # Slide 7: Step 4 - Network Effect
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 4: Network Effect = Selected Referral"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key insight: Network effect IS the referral quality!"
    p = tf.add_paragraph()
    p.text = "G = Selected Referral Quality"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p = tf.add_paragraph()
    p.text = "No mapping needed - they are the same!"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "If selected referral is 'Excellent' → G = Excellent"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "If selected referral is 'Very Bad' → G = Very Bad"
    p.font.size = Pt(18)
    
    # Slide 8: Step 5 - Publication Count
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 5: Publication Count Signal"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Publication count follows Poisson distribution:"
    p = tf.add_paragraph()
    p.text = "• Good candidates: Average 5 papers (λ_good = 5.0)"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "• Bad candidates: Average 2 papers (λ_bad = 2.0)"
    p.font.size = Pt(18)
    p.font.color.rgb = danger_color
    
    # Add visualization
    fig = create_poisson_visualization()
    img_path = 'temp_poisson_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(3), width=Inches(9), height=Inches(4))
    
    # Slide 9: Bayes' Theorem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Step 6: Combine Signals with Bayes' Theorem"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Bayes' Theorem combines prior and signals:"
    p = tf.add_paragraph()
    p.text = "Pr(F = good | G, N) ="
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = "    Pr(G, N | good) × Pr(good)"
    p.font.size = Pt(18)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "    ─────────────────────────────"
    p.font.size = Pt(18)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "         Pr(G, N)"
    p.font.size = Pt(18)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Think of it as: Prior × Evidence = Updated Belief"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    # Slide 10: The Probability Matrix
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Result: Probability Matrix"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "For each combination of signals, we compute the probability:"
    p = tf.add_paragraph()
    p.text = "Rows: Network effect G (5 levels: excellent → very bad)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "Columns: Number of papers N (1-10)"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "Values: Pr(F = good | G, N)"
    p.font.size = Pt(18)
    p.font.bold = True
    
    # Add heatmap
    matrix = create_probability_matrix()
    fig = create_heatmap_figure(matrix)
    img_path = 'temp_heatmap_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.3), Inches(3.5), width=Inches(9.4), height=Inches(3.5))
    
    # Slide 11: Example Scenarios
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Real Examples: Four Candidates"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    # Add visualization
    fig = create_example_scenarios()
    img_path = 'temp_examples_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5.5))
    
    # Slide 12: Key Insights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Insights"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🎯 Trust overrides quality:"
    p = tf.add_paragraph()
    p.text = "   Trusted referral is used even if untrusted is better"
    p.level = 1
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "🛡️ Conservative when both trusted:"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "   When both referrers trusted, use the WORST referral"
    p.level = 1
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "📊 Network effect = Referral quality:"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "   No separate mapping - they are the same!"
    p.level = 1
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✅ Excellent referral + many papers = Strong signal"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "⚠️ Very bad referral + few papers = Weak signal"
    p.font.size = Pt(18)
    p.font.color.rgb = danger_color
    
    # Slide 13: Payoffs Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Beyond Probabilities: Payoffs"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Hiring decisions involve costs and benefits:"
    p = tf.add_paragraph()
    p.text = "💰 B = Payoff from hiring a GOOD candidate (e.g., 10)"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "💰 b = Payoff from hiring a BAD candidate (e.g., -2)"
    p.font.size = Pt(18)
    p.font.color.rgb = danger_color
    p = tf.add_paragraph()
    p.text = "💰 w = Applicant payoff if hired (e.g., 5)"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(52, 152, 219)
    p = tf.add_paragraph()
    p.text = "Key insight: Bad hires can have negative payoffs!"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    # Add visualization
    fig = create_payoff_visualization()
    img_path = 'temp_payoffs_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(3.5), width=Inches(9), height=Inches(3))
    
    # Slide 14: Expected Payoff
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Expected Payoff Calculation"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Expected Payoff = Weighted average of possible outcomes:"
    p = tf.add_paragraph()
    p.text = "E[Payoff] = B × Pr(Good) + b × Pr(Bad)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p = tf.add_paragraph()
    p.text = "Or equivalently:"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "E[Payoff] = b + (B - b) × Pr(Good)"
    p.font.size = Pt(20)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "Example: If Pr(Good) = 0.6, B = 10, b = -2"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "E[Payoff] = -2 + (10 - (-2)) × 0.6 = 5.2"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_color
    
    # Slide 15: Expected Payoff Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Expected Payoff: Real Examples"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    # Add visualization
    fig = create_expected_payoff_comparison()
    img_path = 'temp_expected_payoff_fof.png'
    if output_dir:
        img_path = os.path.join(output_dir, img_path)
    save_figure_to_image(fig, img_path)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), width=Inches(9), height=Inches(5))
    
    # Slide 16: Decision Rules
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Two Decision Approaches"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Approach 1: Probability-Based"
    p = tf.add_paragraph()
    p.text = "IF Pr(F = good | G, N) > 0.5 → HIRE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "ELSE → REJECT"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = danger_color
    p = tf.add_paragraph()
    p.text = "✓ Simple and intuitive"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✗ Assumes symmetric costs"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Approach 2: Payoff-Based"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = primary_color
    p = tf.add_paragraph()
    p.text = "IF E[Payoff] > 0 → HIRE"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = "ELSE → REJECT"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = danger_color
    p = tf.add_paragraph()
    p.text = "✓ Accounts for different costs/benefits"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ More realistic for actual decisions"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✓ Handles negative payoffs from bad hires"
    p.font.size = Pt(18)
    
    # Slide 17: When Decisions Differ
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "When Do Decisions Differ?"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Example: Candidate with Pr(Good) = 0.55"
    p = tf.add_paragraph()
    p.text = "Probability-based: HIRE (55% > 50%)"
    p.font.size = Pt(18)
    p.font.color.rgb = accent_color
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Payoff-based (B=10, b=-2):"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "E[Payoff] = -2 + 12 × 0.55 = 4.6 > 0 → HIRE"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Payoff-based (B=10, b=-10):"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "E[Payoff] = -10 + 20 × 0.55 = 1 > 0 → HIRE"
    p.font.size = Pt(18)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Payoff-based (B=10, b=-30):"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "E[Payoff] = -30 + 40 × 0.55 = -8 < 0 → REJECT"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = danger_color
    p = tf.add_paragraph()
    p.text = "Key insight: High cost of bad hires makes us more conservative!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = primary_color
    
    # Slide 18: Why This Matters
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why This Model Matters"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "✅ Captures realistic trust-based decision making"
    p = tf.add_paragraph()
    p.text = "✅ Models multiple information sources with trust"
    p = tf.add_paragraph()
    p.text = "✅ Quantifies uncertainty (probabilities, not just yes/no)"
    p = tf.add_paragraph()
    p.text = "✅ Shows how trust can override information quality"
    p = tf.add_paragraph()
    p.text = "✅ Conservative approach when multiple trusted sources"
    p = tf.add_paragraph()
    p.text = "✅ Accounts for costs and benefits (payoff modeling)"
    p = tf.add_paragraph()
    p.text = "✅ Transparent and interpretable"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(10)
    
    # Slide 19: Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1️⃣ Start with prior belief about candidate quality"
    p = tf.add_paragraph()
    p.text = "2️⃣ Observe two referrals (R1, R2) and trust (T1, T2)"
    p = tf.add_paragraph()
    p.text = "3️⃣ Select referral based on trust (trust overrides quality!)"
    p = tf.add_paragraph()
    p.text = "4️⃣ Network effect G = selected referral quality"
    p = tf.add_paragraph()
    p.text = "5️⃣ Use Bayes' theorem to get Pr(F = good | G, N)"
    p = tf.add_paragraph()
    p.text = "6️⃣ Compute expected payoff: E[Payoff] = B × Pr(Good) + b × Pr(Bad)"
    p = tf.add_paragraph()
    p.text = "7️⃣ Make hiring decision (probability-based or payoff-based)"
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(10)
        paragraph.font.bold = True
    
    # Slide 20: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Thank You!"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle = slide.placeholders[1]
    subtitle.text = "Questions?"
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Clean up temporary files
    temp_files = ['temp_prior_fof.png', 'temp_referral_dist_fof.png', 'temp_poisson_fof.png', 
                 'temp_heatmap_fof.png', 'temp_examples_fof.png', 'temp_trust_selection_fof.png',
                 'temp_payoffs_fof.png', 'temp_expected_payoff_fof.png']
    for f in temp_files:
        file_path = f if not output_dir else os.path.join(output_dir, f)
        if os.path.exists(file_path):
            os.remove(file_path)
    
    # Determine output file path
    output_filename = "Referral_Based_Hiring_Model.pptx"
    if output_dir:
        # Create directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, output_filename)
    else:
        output_path = output_filename
    
    # Save the presentation
    prs.save(output_path)
    
    return prs, output_path


if __name__ == "__main__":
    import argparse
    import sys
    
    parser = argparse.ArgumentParser(
        description='Create PowerPoint presentation for Referral-Based (FOF) Hiring Model'
    )
    parser.add_argument(
        '--dir', '-d',
        type=str,
        default=None,
        help='Directory where to save the presentation file (default: current directory)'
    )
    
    args = parser.parse_args()
    
    print("Creating PowerPoint presentation for Referral-Based (FOF) model...")
    prs, output_file = create_presentation(output_dir=args.dir)
    print(f"✅ Presentation saved as '{output_file}'")
    print(f"📊 Total slides: {len(prs.slides)}")
